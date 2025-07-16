from abc import ABC, abstractmethod
from typing import List, Dict, Any
import re
import spacy
from src.utils.logger import get_logger

logger = get_logger(__name__)

class BaseExtractor(ABC):
    """Base class for document extractors"""
    
    def __init__(self):
        self.nlp = spacy.load("en_core_web_sm")
    
    @abstractmethod
    def extract_documents(self, limit: int = 50) -> List[Dict[str, Any]]:
        """Extract documents from the source"""
        pass
    
    def clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        if not text:
            return ""
        
        # Basic cleaning
        text = re.sub(r'\s+', ' ', text).strip().lower()
        
        # NLP processing
        doc = self.nlp(text)
        return ' '.join([token.lemma_ for token in doc if not token.is_stop and not token.is_punct])
    
    def extract_text_from_pdf(self, pdf_content: bytes) -> str:
        """Extract text from PDF content"""
        import pdfplumber
        from io import BytesIO
        
        text = ''
        try:
            with pdfplumber.open(BytesIO(pdf_content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text.replace('\n•', '\n- ') + '\n'
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
        
        return text
    
    def extract_text_from_pptx(self, pptx_content: bytes) -> str:
        """Extract text from PPTX content"""
        from pptx import Presentation
        from io import BytesIO
        
        text = ''
        try:
            presentation = Presentation(BytesIO(pptx_content))
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text.replace('\n•', '\n- ') + '\n'
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {e}")
        
        return text 