from abc import ABC, abstractmethod
from typing import List, Dict, Any, Optional, Tuple
import re
from scipy import ndimage
import spacy
from src.utils.logger import get_logger
from io import BytesIO
import requests
logger = get_logger(__name__)

class BaseExtractor(ABC):
    """Base class for document extractors"""

    def __init__(self):
        super().__init__()
        self.nlp = spacy.load("en_core_web_sm")

    @abstractmethod
    def extract_documents(self, limit: int = 50) -> List[Dict[str, Any]]:
        """Extract documents from the source"""
        pass


    def get_chunks(self, text, chunk_size):
        if not text: return []
        start = 0
        end=chunk_size
        logger.info(f"Chunk size is {chunk_size} and num words are {len(text)}")
        if end >= len(text):
            return [' '.join(text)]
        chunks = []
        while end < len(text):
            chunks.append(' '.join(text[start:end]))
            start=end
            end += chunk_size
        return chunks

    def section_get_chunks(self, text, chunk_size):
        from langchain_experimental.text_splitter import SemanticChunker
        from langchain.embeddings import HuggingFaceEmbeddings
        from langchain.schema import Document

        embedding = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L6-v2",
            model_kwargs={"device": "cpu"}  # or "cuda"
        )
        text_splitter = SemanticChunker(embedding)
        text = Document(page_content=' '.join(text))
        try:
            chunks = text_splitter.split_documents([text])
            print(len(chunks))
        except Exception as e:
            logger.error(f"Error in splitting {e}")
        return chunks

    def _download_file(self, link, session, max_file_bytes) -> Optional[Tuple[bytes, Dict[str, str]]]:
        """Download a file with size check; returns (content, headers) or None if skipped/failure."""
        if not link:
            logger.error(f"Invalid link: {link}")
            return None, None
        try:
            assert not isinstance(session, str), type(session)
            # First try HEAD to get size
            try:
                head = session.head(link, allow_redirects=True)
                if head.status_code in (401, 403):
                    logger.warning(f"Unauthorized to download (HEAD): {link}")
                    raise requests.RequestException(f"HEAD returned {head.status_code}")
                if head.status_code == 404:
                    logger.info(f"File not found (404): {link}")
                    return None, None
                if head.status_code >= 400:
                    raise requests.RequestException(f"HEAD returned {head.status_code}")

                size_hdr = head.headers.get("content-length", None)
                size = int(size_hdr) if (size_hdr and size_hdr.isdigit()) else 0
            except requests.RequestException:
                logger.info(f"HEAD failed, falling back to GET for {link}")
                size = 0  # Unknown size, proceed to GET

            # GET request with browser User-Agent
            headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
            resp = session.get(link, headers=headers, stream=True)
            if resp.status_code in (401, 403):
                logger.warning(f"Unauthorized to download (GET): {link}")
                return None, None
            if resp.status_code == 404:
                logger.info(f"File not found (404): {link}")
                return None, None
            if resp.status_code >= 400:
                logger.warning(f"GET failed ({resp.status_code}) for {link}")
                return None, None

            # Use content-length if HEAD succeeded, else unknown
            if size and size > max_file_bytes:
                logger.info(f"Skipping large file ({size} bytes > {max_file_bytes}): {link}")
                return None, None

            # Read content (cap if size unknown)
            if size and size <= max_file_bytes:
                content = resp.content
            else:
                content = resp.raw.read(max_file_bytes + 1)
                if len(content) > max_file_bytes:
                    logger.info(f"Skipping file exceeding max size while streaming: {link}")
                    return None, None

            return content, resp.headers

        except requests.RequestException as e:
            logger.warning(f"Error downloading {link}: {e}")
            return None, None

    def classify_pdf_visual_robust(self, pdf_path):
        """Robust visual classification that handles colored slides"""

        from pdf2image import convert_from_bytes

        # Convert first page ONLY
        try:
            images = convert_from_bytes(pdf_path, first_page=1, last_page=1, dpi=100)
        except:
            logger.error("COuldn't convert from bytes to images")
        width,height = images[0].size
        ratio = width/height
        if ratio > 1.25 and ratio < 1.40 or ratio > 1.7 and ratio < 1.85:
            return 'slides'
        else:
            return 'document'



    def get_raw_text(self, content_type, content):
        raw_text = ''
        if 'application/pdf' in content_type:
            logger.warning("Getting from pdf")
            raw_text, document_type = self.extract_text_from_pdf(content)


            return raw_text, document_type

        elif ('presentationml.presentation' in content_type
                or 'application/vnd.ms-powerpoint' in content_type):
            raw_text,document_type = self.extract_text_from_pptx(content)
        elif ('wordprocessingml.document' in content_type
                or 'application/msword' in content_type
                or 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' in content_type):
            raw_text,document_type = self.extract_text_from_docx(content)
        elif content_type.startswith('text/'):
            raw_text,document_type = self.extract_text_from_txt(content)
        elif ('spreadsheetml.sheet' in content_type
                or 'application/vnd.ms-excel' in content_type
                or 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' in content_type):
            raw_text,document_type = self.extract_text_from_xlsx(content)
        elif content_type in ('image/jpeg', 'image/png', 'image/jpg'):
            raw_text,document_type = self.extract_text_from_image(content)
        else:
            logger.warning(f"Unsupported content type ({content_type})")
        return raw_text, document_type

    def clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        if not text:
            return ""

        # Basic cleaning
        text = re.sub(r'\s+', ' ', text).strip().lower()

        # NLP processing
        doc = self.nlp(text)
        return ' '.join([token.lemma_ for token in doc if not token.is_stop and not token.is_punct])

    def extract_text_from_pdf(self, pdf_content: bytes) -> str: #docdb 126 failing check y  ERROR - Error extracting text from PDF: unsupported operand type(s) for *: 'PSLiteral' and 'float'
        """Extract text from PDF content"""
        import pdfplumber

        text = ''
        try:
            with pdfplumber.open(BytesIO(pdf_content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text.replace('\n•', '\n- ') + '\n'
            document_type = self.classify_pdf_visual_robust(pdf_content)
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")

        return text, document_type

    def extract_text_from_pptx(self, pptx_content: bytes) -> str:
        """Extract text from PPTX content"""
        from pptx import Presentation

        text = ''
        try:
            presentation = Presentation(BytesIO(pptx_content))
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text.replace('\n•', '\n- ') + '\n'
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {e}")
        return text, 'slides'

    def extract_text_from_docx(self, docx_content: bytes) -> str:
        """Extract text from DOCX/DOC content"""
        from docx import Document

        text = ''
        try:
            doc = Document(BytesIO(docx_content))
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + '\n'
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {e}")
        return text, 'document'

    def extract_text_from_txt(self, content: bytes) -> str:
        try:
            return content.decode("utf-8", errors="ignore"), 'document'
        except Exception as e:
            logger.error(f"Error extracting text from TXT: {e}")
        return "", ''

    def extract_text_from_xlsx(self, xlsx_content: bytes) -> str:
        """Extract text from XLSX/XLS spreadsheet"""
        from openpyxl import load_workbook

        text = ''
        try:
            workbook = load_workbook(filename=BytesIO(xlsx_content), data_only=True)
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text += ' '.join(row_text) + '\n'
        except Exception as e:
            logger.error(f"Error extracting text from XLSX: {e}")
        return text, 'slides'

    def extract_text_from_image(self, image_content: bytes) -> str:
        """Extract text from images (JPEG/PNG) using OCR"""
        return ''
        #follow up have this be an llm call to interpret the image in the form of text
        from PIL import Image
        import pytesseract

        text = ''
        try:
            image = Image.open(BytesIO(image_content))
            text = pytesseract.image_to_string(image)
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}")
        return text